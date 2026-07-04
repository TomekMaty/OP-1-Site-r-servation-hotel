"""
Génère diaporama_soutenance.pptx — Maison Saclay — Projet OP-1
10 slides format 16:9
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette ─────────────────────────────────────────────────────────────────
INDIGO     = RGBColor(0x4F, 0x46, 0xE5)   # violet principal
INDIGO_D   = RGBColor(0x43, 0x38, 0xCA)   # violet foncé
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1A, 0x1A, 0x2E)
GRAY       = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG   = RGBColor(0xF8, 0xFA, 0xFF)
BORDER_CLR = RGBColor(0xE0, 0xE7, 0xFF)
ACCENT     = RGBColor(0x10, 0xB9, 0x81)   # vert success
RED        = RGBColor(0xEF, 0x44, 0x44)
LABEL_CLR  = RGBColor(0x9C, 0xA3, 0xAF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

out_path = os.path.join(os.path.dirname(__file__), "diaporama_soutenance.pptx")

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # complètement vide

# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             bold=False, italic=False,
             size=18, color=DARK,
             align=PP_ALIGN.LEFT, wrap=True, valign=None):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if valign:
        from pptx.enum.text import MSO_ANCHOR
        tf.auto_size = None
        tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = Pt(size)
    run.font.color.rgb = color
    return txBox

def add_textbox_multi(slide, lines, l, t, w, h, size=11, color=DARK, bold_first=False, line_spacing=None):
    """lines = list of (text, bold, italic, size_override)"""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, italic, sz = item, False, False, size
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            italic = item[2] if len(item) > 2 else False
            sz = item[3] if len(item) > 3 else size
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.italic = italic
        run.font.size = Pt(sz)
        run.font.color.rgb = color
        if line_spacing:
            from pptx.util import Pt as Pt2
            p.line_spacing = Pt2(line_spacing)
    return txBox

def accent_bar(slide, t=Inches(0), h=Pt(5)):
    """Bande couleur horizontale en haut de la slide."""
    from pptx.dml.color import RGBColor
    shape = slide.shapes.add_shape(1, 0, int(t), int(SLIDE_W), int(h))
    shape.fill.gradient()
    shape.fill.gradient_stops[0].color.rgb = INDIGO
    shape.fill.gradient_stops[1].color.rgb = INDIGO_D
    shape.line.fill.background()
    return shape

def slide_footer(slide, num, total=10):
    # ligne séparatrice
    line = slide.shapes.add_shape(1,
        Inches(0.5), Inches(7.1),
        Inches(12.33), Pt(1))
    line.fill.solid(); line.fill.fore_color.rgb = BORDER_CLR
    line.line.fill.background()
    add_text(slide, "Maison Saclay — Projet OP-1",
             Inches(0.5), Inches(7.15), Inches(6), Pt(18),
             size=8, color=LABEL_CLR)
    add_text(slide, f"Slide {num} / {total}",
             Inches(6.83), Inches(7.15), Inches(6), Pt(18),
             size=8, color=LABEL_CLR, align=PP_ALIGN.RIGHT)

def box_card(slide, l, t, w, h, label, body, label_color=INDIGO):
    add_rect(slide, l, t, w, h, fill=LIGHT_BG, line=BORDER_CLR, line_w=Pt(0.5))
    add_text(slide, label, l+Inches(0.12), t+Inches(0.08), w-Inches(0.24), Pt(14),
             size=7.5, color=label_color, bold=True)
    add_text(slide, body, l+Inches(0.12), t+Inches(0.27), w-Inches(0.24), h-Inches(0.35),
             size=8.5, color=RGBColor(0x37, 0x41, 0x51), wrap=True)

def oral_box(slide, text, t=Inches(6.35)):
    add_rect(slide, Inches(0.5), t, Inches(12.33), Inches(0.72),
             fill=RGBColor(0xF5, 0xF3, 0xFF), line=None)
    # barre gauche
    add_rect(slide, Inches(0.5), t, Pt(4), Inches(0.72), fill=INDIGO)
    add_text(slide, "À dire à l'oral",
             Inches(0.7), t+Inches(0.04), Inches(3), Pt(12),
             size=7.5, color=INDIGO, bold=True)
    add_text(slide, text,
             Inches(0.7), t+Inches(0.22), Inches(11.93), Inches(0.45),
             size=9, color=RGBColor(0x4B, 0x55, 0x63), italic=True, wrap=True)

def slide_header(slide, tag, title, subtitle=None):
    accent_bar(slide)
    add_text(slide, tag,
             Inches(0.5), Inches(0.18), Inches(10), Pt(14),
             size=8, color=INDIGO, bold=True)
    add_text(slide, title,
             Inches(0.5), Inches(0.36), Inches(12.33), Inches(0.8),
             size=26, color=DARK, bold=True)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(1.0), Inches(12.33), Pt(22),
                 size=11, color=GRAY)

def bullet_points(slide, items, l, t, w, h, size=10.5):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        # puce
        run0 = p.add_run()
        run0.text = "●  "
        run0.font.size = Pt(6)
        run0.font.color.rgb = INDIGO
        run0.font.bold = True
        # texte (avec gras pour la partie avant " : ")
        if " : " in item:
            parts = item.split(" : ", 1)
            r1 = p.add_run(); r1.text = parts[0] + " : "
            r1.font.bold = True; r1.font.size = Pt(size); r1.font.color.rgb = DARK
            r2 = p.add_run(); r2.text = parts[1]
            r2.font.bold = False; r2.font.size = Pt(size); r2.font.color.rgb = DARK
        else:
            run = p.add_run(); run.text = item
            run.font.size = Pt(size); run.font.color.rgb = DARK

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITRE / COUVERTURE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)

# Fond foncé gauche
add_rect(sl, 0, 0, Inches(6.2), SLIDE_H, fill=DARK)
# Fond clair droit
add_rect(sl, Inches(6.2), 0, Inches(7.13), SLIDE_H, fill=RGBColor(0xF8, 0xFA, 0xFF))

# Bande violet en haut à gauche
add_rect(sl, 0, 0, Inches(6.2), Pt(5), fill=INDIGO)

# Logo / monogramme
add_rect(sl, Inches(0.6), Inches(1.2), Inches(1.1), Inches(1.1), fill=INDIGO)
add_text(sl, "MS", Inches(0.6), Inches(1.2), Inches(1.1), Inches(1.1),
         size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Titre principal
add_text(sl, "Maison Saclay",
         Inches(0.55), Inches(2.55), Inches(5.4), Inches(0.8),
         size=38, color=WHITE, bold=True)
add_text(sl, "Hôtel 5 étoiles",
         Inches(0.55), Inches(3.25), Inches(5.4), Inches(0.6),
         size=24, color=INDIGO, bold=True)
add_text(sl, "Projet OP-1 — Licence Pro MRT — 2025–2026",
         Inches(0.55), Inches(3.85), Inches(5.4), Inches(0.5),
         size=11, color=RGBColor(0x9C, 0xA3, 0xAF))

# Tags stack côté gauche
tags = ["React + TypeScript", "PHP 8.2 API REST", "MySQL + Docker", "FreeRADIUS + EAP330", "Asterisk 16"]
for i, tag in enumerate(tags):
    ty = Inches(4.65) + i * Inches(0.38)
    add_rect(sl, Inches(0.55), ty, Inches(2.6), Inches(0.31), fill=RGBColor(0x2D, 0x2D, 0x4E), line=INDIGO, line_w=Pt(0.5))
    add_text(sl, tag, Inches(0.65), ty+Pt(3), Inches(2.4), Inches(0.28),
             size=9.5, color=RGBColor(0xC7, 0xD2, 0xFE), bold=True)

# Stats côté droit
stats = [("8", "Services Docker"), ("10", "Tables MySQL"), ("9", "Chambres physiques"), ("RADIUS", "Auth WiFi")]
for i, (val, lbl) in enumerate(stats):
    c = i % 2; r = i // 2
    sx = Inches(6.6) + c * Inches(3.2)
    sy = Inches(1.2) + r * Inches(2.3)
    add_rect(sl, sx, sy, Inches(2.8), Inches(1.8), fill=WHITE, line=BORDER_CLR, line_w=Pt(0.5))
    add_text(sl, val, sx, sy+Inches(0.3), Inches(2.8), Inches(0.9),
             size=36, color=INDIGO, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, lbl, sx, sy+Inches(1.1), Inches(2.8), Pt(22),
             size=9, color=GRAY, align=PP_ALIGN.CENTER)

slide_footer(sl, 1)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — CONTEXTE & OBJECTIF
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "CONTEXTE & OBJECTIF", "Pourquoi ce projet ?")

bullet_points(sl, [
    "Contexte : un hôtel a besoin d'outils numériques pour gérer clients, chambres, commandes et réseau WiFi",
    "Objectif : concevoir et déployer une solution complète, fonctionnelle et démontrable sur le réseau de l'école",
    "Contrainte : tout tourne sur un seul PC via Docker — aucune installation manuelle de PHP, MySQL ou FreeRADIUS",
    "Sécurité : le WiFi n'est accessible qu'après authentification par numéro de chambre (protocole RADIUS)",
], Inches(0.5), Inches(1.3), Inches(12.33), Inches(1.9))

# 2 boxes
box_card(sl, Inches(0.5), Inches(3.35), Inches(6.0), Inches(1.45),
         "CÔTÉ CLIENT",
         "Réserver en ligne · Commander un repas · Accéder au WiFi avec son numéro de chambre · Voir son espace personnel")
box_card(sl, Inches(6.7), Inches(3.35), Inches(6.13), Inches(1.45),
         "CÔTÉ ADMIN",
         "Voir les réservations · Gérer les commandes · Surveiller le WiFi · Recevoir des alertes Telegram en temps réel")

oral_box(sl, '"L\'objectif était de construire un système réaliste, pas juste une maquette. Chaque fonctionnalité fonctionne vraiment : la réservation crée une entrée en base, le WiFi authentifie via RADIUS, les notifications arrivent sur Telegram."')

slide_footer(sl, 2)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ARCHITECTURE DOCKER
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "INFRASTRUCTURE", "Architecture Docker",
             "8 services orchestrés — 1 seule commande : docker compose up -d")

cards = [
    ("nginx",         "Sert le frontend React · Redirige /api/* vers PHP-FPM"),
    ("php 8.2-fpm",   "API REST complète — toute la logique métier"),
    ("mysql 8.0",     "Base de données · Volume persistant (données survivent aux redémarrages)"),
    ("phpmyadmin",    "Interface web DB · Accessible uniquement depuis localhost (sécurisé)"),
    ("freeradius",    "Authentification WiFi RADIUS · Chambres 101–603 configurées"),
    ("asterisk 16",   "Téléphonie IP · 10 extensions SIP (1001-1009 + réception 1099)"),
    ("worker",        "Worker Telegram · Répond aux clics sur les boutons de notification"),
    ("hotel-net",     "Réseau Docker interne · Services communiquent par nom de container"),
]
cols = 4
for i, (label, body) in enumerate(cards):
    c = i % cols; r = i // cols
    bx = Inches(0.5) + c * Inches(3.08)
    by = Inches(1.6) + r * Inches(1.5)
    box_card(sl, bx, by, Inches(2.93), Inches(1.38), label, body)

oral_box(sl, '"Docker permet de déployer l\'ensemble du projet en une seule commande. Chaque service est isolé dans son container. Les données MySQL sont dans un volume persistant — elles survivent aux redémarrages et aux mises à jour."')

slide_footer(sl, 3)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SITE PUBLIC
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "FRONTEND", "Site public hôtel", "SPA React — Single Page Application")

bullet_points(sl, [
    "Page d'accueil : présentation de l'hôtel, navigation vers les sections",
    "Catalogue des chambres : 3 types — Deluxe (38 m²), Suite (72 m²), Penthouse (210 m²)",
    "Formulaire de réservation : choix des dates, nombre de voyageurs, calcul automatique du prix",
    "Page de confirmation : numéro de chambre physique attribuée + extension téléphonique",
    "Room service : commander un repas depuis sa chambre, livraison suivie en temps réel",
], Inches(0.5), Inches(1.35), Inches(12.33), Inches(2.0), size=11)

# Tags
tags_tech = ["React 18", "TypeScript", "Tailwind CSS", "React Router", "Vite (build)", "Framer Motion"]
tx = Inches(0.5)
for tag in tags_tech:
    add_rect(sl, tx, Inches(3.6), Inches(1.55), Inches(0.32), fill=RGBColor(0xED, 0xE9, 0xFE), line=RGBColor(0xC4, 0xB5, 0xFD), line_w=Pt(0.4))
    add_text(sl, tag, tx+Pt(4), Inches(3.62), Inches(1.47), Inches(0.28), size=8, color=RGBColor(0x43, 0x38, 0xCA), bold=True)
    tx += Inches(1.65)

box_card(sl, Inches(0.5), Inches(4.15), Inches(12.33), Inches(1.1),
         "COMMUNICATION FRONTEND → BACKEND",
         "Centralisée dans api.ts — ajoute automatiquement le token JWT dans chaque requête HTTP. Le backend répond toujours en JSON avec success: true/false.")

oral_box(sl, '"Le frontend est une SPA React compilée avec Vite. Il communique avec le backend PHP uniquement via des appels API JSON. Nginx sert les fichiers statiques du build et redirige les requêtes /api vers PHP."')

slide_footer(sl, 4)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — RÉSERVATION ET BASE DE DONNÉES
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "BACKEND + MYSQL", "Réservation et base de données")

# Flow
flow_steps = [
    "1. Client choisit les dates",
    "→",
    "2. POST /api/bookings + JWT",
    "→",
    "3. Dispo SQL",
    "→",
    "4. Chambre physique attribuée",
    "→",
    "5. INSERT MySQL",
    "→",
    "6. Telegram",
]
add_rect(sl, Inches(0.5), Inches(1.35), Inches(12.33), Inches(0.6),
         fill=RGBColor(0xF8, 0xFA, 0xFF), line=BORDER_CLR, line_w=Pt(0.5))
tx = Inches(0.6)
for step in flow_steps:
    is_arrow = step == "→"
    w = Pt(18) if is_arrow else Inches(1.88)
    clr = INDIGO if is_arrow else DARK
    sz = 16 if is_arrow else 7.5
    bold = is_arrow
    if not is_arrow:
        add_rect(sl, tx, Inches(1.4), w, Inches(0.5),
                 fill=WHITE, line=RGBColor(0xC7, 0xD2, 0xFE), line_w=Pt(0.4))
    add_text(sl, step, tx, Inches(1.45), w, Inches(0.45),
             size=sz, color=clr, bold=bold, align=PP_ALIGN.CENTER)
    tx += w + Pt(4)

bullet_points(sl, [
    "Attribution automatique : requête SQL avec détection de chevauchement de dates — aucune double-réservation possible",
    "Prix calculé : nombre de nuits × tarif du type de chambre",
    "Notification Telegram : l'admin reçoit un message avec boutons Confirmer / Annuler interactifs",
    "Chambre physique : distinction entre le type (catalogue) et la chambre réelle attribuée (101, 201, 601…)",
], Inches(0.5), Inches(2.1), Inches(12.33), Inches(1.8), size=10.5)

# Stats row
stats5 = [("3", "Types"), ("9", "Chambres physiques"), ("16", "Réservations"), ("11", "Commandes"), ("10", "Tables MySQL")]
for i, (val, lbl) in enumerate(stats5):
    sx = Inches(0.5) + i * Inches(2.42)
    add_rect(sl, sx, Inches(4.05), Inches(2.28), Inches(1.0),
             fill=LIGHT_BG, line=BORDER_CLR, line_w=Pt(0.5))
    add_text(sl, val, sx, Inches(4.12), Inches(2.28), Inches(0.5),
             size=26, color=INDIGO, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, lbl, sx, Inches(4.6), Inches(2.28), Pt(18),
             size=8.5, color=GRAY, align=PP_ALIGN.CENTER)

oral_box(sl, '"La réservation déclenche une série d\'opérations : vérification SQL avec chevauchement de dates, attribution automatique de chambre physique, calcul du prix et notification Telegram avec boutons interactifs pour l\'admin."')

slide_footer(sl, 5)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DASHBOARD ADMIN
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "ADMIN REACT", "Dashboard administrateur",
             "Route /admin — protégée par mot de passe + token JWT")

# 6 cards en 3+3
tab_cards = [
    ("RÉSERVATIONS", "Liste complète · Numéro chambre attribué · Changer le statut"),
    ("COMMANDES", "Room service · Articles · Statuts : en attente / préparation / livré"),
    ("CHAMBRES", "Vue temps réel : libre, occupée, arrivée aujourd'hui, départ…"),
    ("WIFI", "Sessions RADIUS réelles + activité réseau de démonstration"),
    ("STATISTIQUES", "Revenus, taux d'occupation, activité du jour"),
    ("CHARGEMENT", "Toutes les données en parallèle via Promise.all() au démarrage"),
]
for i, (lbl, body) in enumerate(tab_cards):
    c = i % 3; r = i // 3
    bx = Inches(0.5) + c * Inches(4.15)
    by = Inches(1.55) + r * Inches(1.45)
    box_card(sl, bx, by, Inches(4.0), Inches(1.32), lbl, body)

# 2 auth boxes
box_card(sl, Inches(0.5), Inches(4.6), Inches(6.0), Inches(1.15),
         "AUTHENTIFICATION ADMIN",
         "Mot de passe unique (.env) → JWT signé HS256 · Valable 7 jours · Stocké en sessionStorage")
box_card(sl, Inches(6.7), Inches(4.6), Inches(6.13), Inches(1.15),
         "PROTECTION API",
         "AdminMiddleware vérifie le JWT à chaque requête → HTTP 401 si absent ou invalide")

oral_box(sl, '"L\'admin est une SPA React sécurisée par JWT. Toutes les données sont chargées en parallèle au démarrage. Chaque route API admin retourne HTTP 401 si le token est absent ou expiré — il n\'y a aucun accès sans authentification."')

slide_footer(sl, 6)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PORTAIL CAPTIF WIFI
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "WIFI — TP-LINK EAP330", "Portail captif WiFi",
             "Authentification des clients par numéro de chambre")

# Flow
flow7 = [
    "Téléphone → OP1_Wifi",
    "→",
    "Portail EAP330",
    "→",
    "Login : 201 / 201",
    "→",
    "RADIUS :18120",
    "→",
    "FreeRADIUS Docker",
    "→",
    "Access-Accept → Internet OK",
]
add_rect(sl, Inches(0.5), Inches(1.45), Inches(12.33), Inches(0.65),
         fill=RGBColor(0xF8, 0xFA, 0xFF), line=BORDER_CLR, line_w=Pt(0.5))
tx = Inches(0.62)
for step in flow7:
    is_arrow = step == "→"
    w = Pt(16) if is_arrow else Inches(1.72)
    clr = INDIGO if is_arrow else DARK
    sz = 14 if is_arrow else 7.5
    if not is_arrow:
        add_rect(sl, tx, Inches(1.5), w, Inches(0.52),
                 fill=WHITE, line=RGBColor(0xC7, 0xD2, 0xFE), line_w=Pt(0.4))
    add_text(sl, step, tx, Inches(1.55), w, Inches(0.45),
             size=sz, color=clr, bold=is_arrow, align=PP_ALIGN.CENTER)
    tx += w + Pt(3)

bullet_points(sl, [
    "Identifiant = Mot de passe = Numéro de chambre (101 à 603)",
    "Borne EAP330 en mode \"Local Web Portal\" + RADIUS externe",
    "Session autorisée 1 heure — puis re-authentification requise",
    "FreeRADIUS répond Access-Accept (OK) ou Access-Reject (refusé)",
    "RADIUS = protocole standard en hôtellerie, entreprises et universités",
], Inches(0.5), Inches(2.28), Inches(12.33), Inches(2.0), size=11)

# Badge démo
add_rect(sl, Inches(0.5), Inches(4.5), Inches(5.8), Inches(0.38),
         fill=RGBColor(0xEC, 0xFD, 0xF5), line=RGBColor(0x6E, 0xE7, 0xB7), line_w=Pt(0.5))
add_text(sl, "📱  Démo live : connexion WiFi depuis un téléphone — chambre 201",
         Inches(0.65), Inches(4.54), Inches(5.6), Inches(0.32),
         size=9, color=RGBColor(0x06, 0x5F, 0x46), bold=True)

oral_box(sl, '"Le portail captif utilise le protocole RADIUS, standard dans les réseaux hôteliers. La borne TP-Link délègue l\'authentification à FreeRADIUS dans Docker. Le client entre son numéro de chambre et accède à Internet."')

slide_footer(sl, 7)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FREERADIUS + ADMIN WIFI
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "FREERADIUS + REMONTÉE ADMIN", "Sessions WiFi dans l'admin")

bullet_points(sl, [
    "FreeRADIUS dans Docker : 10 chambres (101–603) configurées, secret partagé avec la borne",
    "Relay UDP PowerShell : contourne la limitation Docker Desktop Windows — écoute :18120, retransmet vers FreeRADIUS :1812",
    "Script sync_radius_sessions.sh : parse les logs FreeRADIUS, insère/met à jour MySQL (idempotent — sans doublons)",
    "Onglet WiFi de l'admin : appareils connectés, chambres, adresses MAC, Access-Accept/Reject",
    "Activité réseau : section de démonstration honnêtement labelisée (domaines DNS consultés)",
], Inches(0.5), Inches(1.35), Inches(12.33), Inches(2.2), size=11)

box_card(sl, Inches(0.5), Inches(3.75), Inches(6.0), Inches(1.15),
         "DONNÉES WIFI EN BASE (RÉELLES)",
         "4 sessions · 3 Access-Accept · 1 Reject · Adresses MAC réelles d'appareils testés")
box_card(sl, Inches(6.7), Inches(3.75), Inches(6.13), Inches(1.15),
         "ACTIVITÉ RÉSEAU (DÉMO)",
         "10 lignes DNS · Chambres 101 et 201 · Clairement marquées « démonstration » dans l'interface")

add_rect(sl, Inches(0.5), Inches(5.1), Inches(7.5), Inches(0.38),
         fill=RGBColor(0xEC, 0xFD, 0xF5), line=RGBColor(0x6E, 0xE7, 0xB7), line_w=Pt(0.5))
add_text(sl, "💻  Démo : bash sync_radius_sessions.sh → admin WiFi mis à jour",
         Inches(0.65), Inches(5.14), Inches(7.3), Inches(0.32),
         size=9, color=RGBColor(0x06, 0x5F, 0x46), bold=True)

oral_box(sl, '"Un script bash lit les logs FreeRADIUS et les transforme en données SQL. L\'admin voit en temps réel quels appareils sont connectés, depuis quelle chambre. Le script peut être relancé sans créer de doublons grâce à ON DUPLICATE KEY UPDATE."')

slide_footer(sl, 8)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — SÉCURITÉ + DIFFICULTÉS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "SÉCURITÉ & PROBLÈMES RÉSOLUS", "Difficultés et solutions")

# Security cards
sec_cards = [
    ("JWT HS256",          "Tokens signés et vérifiés à chaque requête — impossible à forger sans le secret"),
    ("Mots de passe bcrypt","Jamais stockés en clair — hash bcrypt irréversible en base MySQL"),
    ("Requêtes PDO préparées","Paramètres bindés — protection totale contre les injections SQL"),
    ("Ports restreints",   "phpMyAdmin + MySQL limités à 127.0.0.1 — inaccessibles depuis le réseau WiFi"),
]
for i, (lbl, body) in enumerate(sec_cards):
    c = i % 2; r = i // 2
    bx = Inches(0.5) + c * Inches(6.22)
    by = Inches(1.55) + r * Inches(1.1)
    box_card(sl, bx, by, Inches(6.08), Inches(1.0), lbl, body)

# ligne séparatrice
add_rect(sl, Inches(0.5), Inches(3.85), Inches(12.33), Pt(1), fill=BORDER_CLR)

add_text(sl, "PROBLÈMES RENCONTRÉS ET SOLUTIONS",
         Inches(0.5), Inches(3.95), Inches(12.33), Pt(14),
         size=7.5, color=LABEL_CLR, bold=True)

# Problème 1 + solution
add_rect(sl, Inches(0.5), Inches(4.2), Inches(5.9), Inches(0.98),
         fill=LIGHT_BG, line=RED, line_w=Pt(1.5))
add_text(sl, "PROBLÈME 1 — Docker + UDP",
         Inches(0.62), Inches(4.24), Inches(5.7), Pt(13), size=7.5, color=RED, bold=True)
add_text(sl, "Docker Desktop Windows ne peut pas recevoir de paquets UDP depuis le réseau local (NAT Docker)",
         Inches(0.62), Inches(4.42), Inches(5.7), Inches(0.7), size=9, color=DARK, wrap=True)

add_rect(sl, Inches(6.6), Inches(4.2), Inches(6.13), Inches(0.98),
         fill=LIGHT_BG, line=ACCENT, line_w=Pt(1.5))
add_text(sl, "SOLUTION",
         Inches(6.72), Inches(4.24), Inches(5.9), Pt(13), size=7.5, color=ACCENT, bold=True)
add_text(sl, "Relay UDP PowerShell : écoute sur 192.168.30.104:18120 et retransmet vers FreeRADIUS sur 127.0.0.1:1812",
         Inches(6.72), Inches(4.42), Inches(5.9), Inches(0.7), size=9, color=DARK, wrap=True)

# Problème 2 + solution
add_rect(sl, Inches(0.5), Inches(5.3), Inches(5.9), Inches(0.88),
         fill=LIGHT_BG, line=RED, line_w=Pt(1.5))
add_text(sl, "PROBLÈME 2 — Double réservation",
         Inches(0.62), Inches(5.34), Inches(5.7), Pt(13), size=7.5, color=RED, bold=True)
add_text(sl, "Deux clients pouvaient réserver la même chambre aux mêmes dates",
         Inches(0.62), Inches(5.52), Inches(5.7), Inches(0.62), size=9, color=DARK, wrap=True)

add_rect(sl, Inches(6.6), Inches(5.3), Inches(6.13), Inches(0.88),
         fill=LIGHT_BG, line=ACCENT, line_w=Pt(1.5))
add_text(sl, "SOLUTION",
         Inches(6.72), Inches(5.34), Inches(5.9), Pt(13), size=7.5, color=ACCENT, bold=True)
add_text(sl, "Requête SQL avec condition de chevauchement (check_in < sortie AND check_out > arrivée) + attribution auto",
         Inches(6.72), Inches(5.52), Inches(5.9), Inches(0.62), size=9, color=DARK, wrap=True)

oral_box(sl, '"La principale difficulté était Docker + RADIUS sur Windows. On a résolu ça avec un relay UDP maison en PowerShell. C\'est une solution non standard, mais stable et documentée."', t=Inches(6.32))

slide_footer(sl, 9)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CONCLUSION
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "BILAN", "Conclusion et améliorations")

add_text(sl, "Ce qui fonctionne — démontrable en live",
         Inches(0.5), Inches(1.35), Inches(12.33), Pt(14),
         size=8, color=LABEL_CLR, bold=True)

bullet_points(sl, [
    "Réservation multi-nuits de bout en bout avec attribution automatique de chambre physique",
    "Room service : commande → Telegram → suivi statut dans l'admin",
    "Portail captif WiFi RADIUS — testé avec de vrais appareils sur le réseau",
    "Dashboard admin complet : réservations, commandes, chambres, WiFi, statistiques",
    "Infrastructure Docker stable — 8 services, IP statique, relay UDP persistant",
], Inches(0.5), Inches(1.62), Inches(12.33), Inches(1.85), size=11)

add_rect(sl, Inches(0.5), Inches(3.6), Inches(12.33), Pt(1), fill=BORDER_CLR)

add_text(sl, "AMÉLIORATIONS POSSIBLES",
         Inches(0.5), Inches(3.72), Inches(12.33), Pt(14),
         size=8, color=LABEL_CLR, bold=True)

improv = [
    ("HTTPS",       "Certificat SSL Let's Encrypt pour chiffrer tous les échanges"),
    ("DNS réel",    "Pi-hole pour remplacer la démo par une vraie activité réseau"),
    ("Téléphonie",  "Connecter des postes SIP physiques aux extensions Asterisk"),
    ("Monitoring",  "Alertes automatiques si un container Docker tombe"),
]
for i, (lbl, body) in enumerate(improv):
    bx = Inches(0.5) + i * Inches(3.1)
    box_card(sl, bx, Inches(4.0), Inches(2.95), Inches(1.2), lbl, body)

oral_box(sl, '"Le projet est fonctionnel et démontrable dans son ensemble. Toutes les fonctionnalités principales marchent en conditions réelles. Avec plus de temps, on ajouterait HTTPS, un vrai DNS local et des postes téléphoniques physiques."', t=Inches(5.38))

add_rect(sl, 0, Inches(6.28), SLIDE_W, Inches(1.22), fill=DARK)
add_text(sl, "Maison Saclay — Hôtel 5 étoiles — Projet OP-1",
         0, Inches(6.48), SLIDE_W, Pt(24),
         size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, "Licence Pro MRT — 2025–2026",
         0, Inches(6.82), SLIDE_W, Pt(20),
         size=10, color=GRAY, align=PP_ALIGN.CENTER)

slide_footer(sl, 10)

# ── Sauvegarde ───────────────────────────────────────────────────────────────
prs.save(out_path)
print(f"✅ Fichier créé : {out_path}")
