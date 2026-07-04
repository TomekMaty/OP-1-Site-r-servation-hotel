"""
Génère diaporama_soutenance_v2.pptx — Maison Saclay
Version éditable : layouts PowerPoint standards, placeholders, tableaux natifs
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree
import copy, os

# ── Palette ──────────────────────────────────────────────────────────────────
INDIGO   = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_D = RGBColor(0x3B, 0x35, 0xC0)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x1A, 0x1A, 0x2E)
GRAY     = RGBColor(0x6B, 0x72, 0x80)
LIGHT    = RGBColor(0xF8, 0xFA, 0xFF)
BORDER   = RGBColor(0xC7, 0xD2, 0xFE)
GREEN    = RGBColor(0x06, 0x5F, 0x46)
RED_C    = RGBColor(0xDC, 0x26, 0x26)

out_path = os.path.join(os.path.dirname(__file__), "diaporama_soutenance_v2.pptx")

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# On utilise les layouts disponibles
LAY_TITLE      = prs.slide_layouts[0]   # Grand titre + sous-titre
LAY_CONTENT    = prs.slide_layouts[1]   # Titre + zone de contenu
LAY_TWO_COL    = prs.slide_layouts[3]   # Titre + 2 colonnes
LAY_BLANK      = prs.slide_layouts[6]   # Vierge (pour slides custom)

# ── Helpers généraux ─────────────────────────────────────────────────────────

def rgb(r,g,b): return RGBColor(r,g,b)

def set_font(run, size=None, bold=None, color=None, italic=None):
    if size:   run.font.size  = Pt(size)
    if bold is not None: run.font.bold = bold
    if color:  run.font.color.rgb = color
    if italic is not None: run.font.italic = italic

def bar_top(slide, color=INDIGO):
    """Bande de couleur fine en haut de la slide — facile à supprimer/modifier."""
    sh = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Pt(6))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.name = "TopBar"

def footer(slide, num):
    """Pied de page : nom du projet + numéro de slide."""
    line = slide.shapes.add_shape(1,
        Inches(0.4), Inches(7.1), Inches(12.5), Pt(1))
    line.fill.solid(); line.fill.fore_color.rgb = rgb(0xE5,0xE7,0xEB)
    line.line.fill.background(); line.name = "FooterLine"

    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.15), Inches(12.5), Pt(20))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "Maison Saclay — Projet OP-1"
    set_font(r1, size=8, color=GRAY)
    r2 = p.add_run(); r2.text = f"         {num} / 10"
    set_font(r2, size=8, color=GRAY)
    p.alignment = PP_ALIGN.LEFT
    tb.name = "Footer"

def tag_box(slide, text, l, t, color=INDIGO, bg=LIGHT):
    """Petite étiquette colorée (tag/badge) — comme les tags technos."""
    w = Inches(1.6); h = Inches(0.32)
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = bg
    sh.line.color.rgb = color; sh.line.width = Pt(0.75)
    tb2 = sh.text_frame
    tb2.word_wrap = False
    p = tb2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    set_font(r, size=9, bold=True, color=color)
    sh.name = f"Tag_{text[:8]}"
    return sh

def card(slide, l, t, w, h, title, body, title_color=INDIGO):
    """Carte avec titre en couleur et corps de texte — entièrement éditable."""
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = LIGHT
    sh.line.color.rgb = BORDER; sh.line.width = Pt(0.75)
    sh.name = f"Card_{title[:8]}"
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.12); tf.margin_top = Pt(6)
    tf.margin_right = Inches(0.12); tf.margin_bottom = Pt(6)

    p1 = tf.paragraphs[0]
    r1 = p1.add_run(); r1.text = title.upper()
    set_font(r1, size=7.5, bold=True, color=title_color)

    p2 = tf.add_paragraph(); p2.space_before = Pt(4)
    r2 = p2.add_run(); r2.text = body
    set_font(r2, size=9, color=DARK)

def oral(slide, text, t=Inches(6.35)):
    """Encadré 'À dire à l'oral' — violet clair, facilement supprimable."""
    bg = slide.shapes.add_shape(1, Inches(0.4), t, Inches(12.5), Inches(0.72))
    bg.fill.solid(); bg.fill.fore_color.rgb = rgb(0xED,0xE9,0xFE)
    bg.line.fill.background(); bg.name = "OralBg"

    accent = slide.shapes.add_shape(1, Inches(0.4), t, Pt(5), Inches(0.72))
    accent.fill.solid(); accent.fill.fore_color.rgb = INDIGO
    accent.line.fill.background(); accent.name = "OralAccent"

    tb = slide.shapes.add_textbox(Inches(0.65), t+Pt(5), Inches(12.2), Inches(0.62))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r0 = p.add_run(); r0.text = "À dire à l'oral  "
    set_font(r0, size=7.5, bold=True, color=INDIGO)
    r1 = p.add_run(); r1.text = text
    set_font(r1, size=9, italic=True, color=rgb(0x4B,0x55,0x63))
    tb.name = "OralText"

def title_area(slide, tag, title, subtitle=None):
    """En-tête standard : tag + titre H1 + sous-titre optionnel."""
    bar_top(slide)
    tb_tag = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(10), Pt(14))
    p = tb_tag.text_frame.paragraphs[0]
    r = p.add_run(); r.text = tag.upper()
    set_font(r, size=8, bold=True, color=INDIGO)
    tb_tag.name = "SlideTag"

    tb_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.75))
    tf = tb_title.text_frame; tf.word_wrap = True
    p2 = tf.paragraphs[0]
    r2 = p2.add_run(); r2.text = title
    set_font(r2, size=28, bold=True, color=DARK)
    tb_title.name = "SlideTitle"

    if subtitle:
        tb_sub = slide.shapes.add_textbox(Inches(0.5), Inches(0.98), Inches(12.3), Pt(22))
        p3 = tb_sub.text_frame.paragraphs[0]
        r3 = p3.add_run(); r3.text = subtitle
        set_font(r3, size=11, color=GRAY)
        tb_sub.name = "SlideSubtitle"

def bullets(slide, items, l, t, w, h, size=11):
    """
    Liste à puces — utilise un textbox standard avec paragraphes séparés.
    Dans PowerPoint : clic sur le bloc → tu peux modifier chaque ligne.
    Gras automatique pour la partie avant ' : '
    """
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    tb.name = "BulletList"
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(4)
        # Puce colorée
        r0 = p.add_run(); r0.text = "●  "
        set_font(r0, size=6, bold=True, color=INDIGO)
        # Texte avec gras sur la partie avant " : "
        if " : " in item:
            a, b = item.split(" : ", 1)
            r1 = p.add_run(); r1.text = a + " : "
            set_font(r1, size=size, bold=True, color=DARK)
            r2 = p.add_run(); r2.text = b
            set_font(r2, size=size, color=DARK)
        else:
            r1 = p.add_run(); r1.text = item
            set_font(r1, size=size, color=DARK)

def add_table(slide, headers, rows, l, t, w, h, header_bg=INDIGO):
    """
    Tableau natif PowerPoint — tu peux cliquer sur chaque cellule et modifier.
    Contrairement aux shapes texte, c'est un VRAI tableau PowerPoint.
    """
    tbl = slide.shapes.add_table(len(rows)+1, len(headers), l, t, w, h).table
    tbl.first_row = True

    # En-têtes
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_bg
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = hdr
        set_font(r, size=9, bold=True, color=WHITE)

    # Données
    for ri, row in enumerate(rows):
        bg = LIGHT if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            set_font(r, size=9, color=DARK)

    return tbl

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COUVERTURE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_BLANK)
bar_top(sl, INDIGO)

# Fond sombre gauche
bg_left = sl.shapes.add_shape(1, 0, 0, Inches(6.4), prs.slide_height)
bg_left.fill.solid(); bg_left.fill.fore_color.rgb = DARK
bg_left.line.fill.background(); bg_left.name = "CoverBgLeft"

# Logo MS
logo = sl.shapes.add_shape(1, Inches(0.6), Inches(1.1), Inches(1.1), Inches(1.1))
logo.fill.solid(); logo.fill.fore_color.rgb = INDIGO
logo.line.fill.background(); logo.name = "LogoBox"
tf_l = logo.text_frame; tf_l.paragraphs[0].alignment = PP_ALIGN.CENTER
rl = tf_l.paragraphs[0].add_run(); rl.text = "MS"
set_font(rl, size=28, bold=True, color=WHITE)

# Titre principal
tb_title = sl.shapes.add_textbox(Inches(0.55), Inches(2.5), Inches(5.5), Inches(1.0))
tf = tb_title.text_frame; tf.word_wrap = True; tb_title.name = "CoverTitle"
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Maison Saclay"
set_font(r, size=40, bold=True, color=WHITE)

tb_sub1 = sl.shapes.add_textbox(Inches(0.55), Inches(3.35), Inches(5.5), Pt(30))
p2 = tb_sub1.text_frame.paragraphs[0]; tb_sub1.name = "CoverSubtitle1"
r2 = p2.add_run(); r2.text = "Hôtel 5 étoiles — Projet OP-1"
set_font(r2, size=16, color=INDIGO)

tb_sub2 = sl.shapes.add_textbox(Inches(0.55), Inches(3.85), Inches(5.5), Pt(22))
p3 = tb_sub2.text_frame.paragraphs[0]; tb_sub2.name = "CoverSubtitle2"
r3 = p3.add_run(); r3.text = "Licence Pro MRT — 2025–2026"
set_font(r3, size=10.5, color=GRAY)

# Tags techno (modifiables individuellement)
tags = ["React + TypeScript", "PHP 8.2 API REST", "MySQL + Docker", "FreeRADIUS + EAP330", "Asterisk 16"]
for i, tg in enumerate(tags):
    tag_box(sl, tg, Inches(0.55), Inches(4.65)+i*Inches(0.38))

# Stats côté droit (fond clair)
stats = [("8", "Services Docker"), ("10", "Tables MySQL"), ("9", "Chambres physiques"), ("RADIUS", "Auth WiFi")]
for i, (val, lbl) in enumerate(stats):
    c = i % 2; r_i = i // 2
    sx = Inches(6.8) + c*Inches(3.1)
    sy = Inches(1.3) + r_i*Inches(2.5)
    sh = sl.shapes.add_shape(1, sx, sy, Inches(2.9), Inches(2.0))
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(0xF0,0xF0,0xFF)
    sh.line.color.rgb = BORDER; sh.line.width = Pt(0.5); sh.name = f"Stat_{val}"
    tf2 = sh.text_frame
    tf2.margin_top = Inches(0.25)
    p_v = tf2.paragraphs[0]; p_v.alignment = PP_ALIGN.CENTER
    r_v = p_v.add_run(); r_v.text = val
    set_font(r_v, size=38, bold=True, color=INDIGO)
    p_l = tf2.add_paragraph(); p_l.alignment = PP_ALIGN.CENTER
    r_l = p_l.add_run(); r_l.text = lbl
    set_font(r_l, size=9, color=GRAY)

footer(sl, 1)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — CONTEXTE & OBJECTIF
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_BLANK)
title_area(sl, "Contexte & objectif", "Pourquoi ce projet ?")

bullets(sl, [
    "Contexte : un hôtel a besoin d'outils numériques pour gérer clients, chambres, commandes et réseau WiFi",
    "Objectif : concevoir et déployer une solution complète, fonctionnelle et démontrable sur le réseau de l'école",
    "Contrainte : tout tourne sur un seul PC via Docker — aucune installation manuelle de PHP, MySQL ou FreeRADIUS",
    "Sécurité : le WiFi n'est accessible qu'après authentification par numéro de chambre (protocole RADIUS)",
], Inches(0.5), Inches(1.35), Inches(12.3), Inches(1.9))

card(sl, Inches(0.5), Inches(3.4), Inches(6.0), Inches(1.5),
     "Côté Client",
     "Réserver en ligne · Commander un repas · Accéder au WiFi avec son numéro de chambre · Voir son espace personnel")
card(sl, Inches(6.75), Inches(3.4), Inches(6.08), Inches(1.5),
     "Côté Admin",
     "Voir les réservations · Gérer les commandes · Surveiller le WiFi · Recevoir des alertes Telegram en temps réel")

oral(sl, '"L\'objectif était de construire un système réaliste, pas juste une maquette. Chaque fonctionnalité fonctionne vraiment : la réservation crée une entrée en base, le WiFi authentifie via RADIUS, les notifications arrivent sur Telegram."')
footer(sl, 2)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ARCHITECTURE DOCKER — TABLEAU NATIF
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_BLANK)
title_area(sl, "Infrastructure", "Architecture Docker",
           "8 services orchestrés — une seule commande : docker compose up -d")

headers = ["Container", "Image", "Rôle"]
rows = [
    ("maison_nginx",       "nginx:alpine",          "Sert le frontend React · redirige /api/* vers PHP"),
    ("maison_php",         "php:8.2-fpm",           "API REST — toute la logique métier"),
    ("maison_mysql",       "mysql:8.0",             "Base de données · volume persistant"),
    ("maison_freeradius",  "freeradius custom",     "Authentification WiFi RADIUS · chambres 101–603"),
    ("maison_worker",      "php:8.2-fpm",           "Worker Telegram · traite les clics sur boutons"),
    ("maison_asterisk",    "andrius/asterisk:16",   "Téléphonie IP · 10 extensions SIP (1001–1099)"),
    ("maison_phpmyadmin",  "phpmyadmin/phpmyadmin", "Interface web DB · localhost:8080"),
    ("maison_portail",     "php:8.2-fpm",           "Ancien portail captif (legacy)"),
]
add_table(sl, headers, rows, Inches(0.4), Inches(1.35), Inches(12.5), Inches(4.85))

oral(sl, '"Docker permet de déployer l\'ensemble du projet en une seule commande. Chaque service est isolé dans son container. Les données MySQL sont dans un volume persistant — elles survivent aux redémarrages."')
footer(sl, 3)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SITE PUBLIC
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_BLANK)
title_area(sl, "Frontend React", "Site public hôtel", "SPA — Single Page Application")

bullets(sl, [
    "Page d'accueil : présentation de l'hôtel, navigation vers les sections",
    "Catalogue des chambres : 3 types — Deluxe (38 m²), Suite (72 m²), Penthouse (210 m²)",
    "Formulaire de réservation : dates, voyageurs, calcul automatique du prix",
    "Page de confirmation : numéro de chambre physique attribuée + extension téléphonique",
    "Room service : commander un repas depuis sa chambre, suivi en temps réel",
], Inches(0.5), Inches(1.35), Inches(7.5), Inches(2.1))

# Tags techno à droite
tech = ["React 18", "TypeScript", "Tailwind CSS", "React Router", "Vite", "Framer Motion"]
for i, tg in enumerate(tech):
    c = i % 2; r_i = i // 2
    tag_box(sl, tg, Inches(8.2)+c*Inches(1.72), Inches(1.45)+r_i*Inches(0.45))

card(sl, Inches(0.5), Inches(3.65), Inches(12.3), Inches(1.05),
     "Communication Frontend → Backend",
     "Centralisée dans api.ts — ajoute automatiquement le token JWT dans chaque requête. "
     "Le backend répond toujours en JSON avec success: true/false.")

oral(sl, '"Le frontend est une SPA React compilée avec Vite. Il communique avec le backend PHP uniquement via des appels API JSON. Nginx sert les fichiers statiques et redirige les requêtes /api vers PHP."')
footer(sl, 4)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — RÉSERVATION & BASE DE DONNÉES
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_BLANK)
title_area(sl, "Backend + MySQL", "Réservation et base de données")

# Flow en forme de cartes enchaînées
flow = ["1. Client choisit les dates", "2. POST /api/bookings + JWT",
        "3. Dispo SQL", "4. Chambre attribuée", "5. INSERT MySQL", "6. Telegram"]
tx = Inches(0.4)
bg_flow = sl.shapes.add_shape(1, Inches(0.4), Inches(1.35), Inches(12.5), Inches(0.65))
bg_flow.fill.solid(); bg_flow.fill.fore_color.rgb = LIGHT
bg_flow.line.color.rgb = BORDER; bg_flow.line.width = Pt(0.5); bg_flow.name = "FlowBg"
for i, step in enumerate(flow):
    w = Inches(1.95)
    sh = sl.shapes.add_shape(1, tx, Inches(1.4), w, Inches(0.54))
    sh.fill.solid(); sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = BORDER; sh.line.width = Pt(0.5); sh.name = f"Flow{i}"
    tf2 = sh.text_frame; tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    r2 = tf2.paragraphs[0].add_run(); r2.text = step
    set_font(r2, size=7.5, color=DARK, bold=True)
    tx += w + Inches(0.12)
    if i < len(flow)-1:
        arr = sl.shapes.add_textbox(tx-Inches(0.16), Inches(1.47), Inches(0.14), Pt(22))
        arr.text_frame.paragraphs[0].add_run().text = "→"
        set_font(arr.text_frame.paragraphs[0].runs[0], size=12, bold=True, color=INDIGO)
        arr.name = f"Arrow{i}"

bullets(sl, [
    "Attribution automatique : SQL avec chevauchement de dates — aucune double-réservation possible",
    "Prix calculé : nombre de nuits × tarif du type de chambre",
    "Notification Telegram : message avec boutons Confirmer / Annuler interactifs",
    "Chambre physique : distinction entre type catalogue et chambre réelle (101, 201, 601…)",
], Inches(0.5), Inches(2.1), Inches(12.3), Inches(1.75))

# Stats
stats5 = [("3", "Types"), ("9", "Chambres"), ("16", "Réservations"), ("11", "Commandes"), ("10", "Tables SQL")]
for i, (val, lbl) in enumerate(stats5):
    sx = Inches(0.4) + i*Inches(2.48)
    sh = sl.shapes.add_shape(1, sx, Inches(4.05), Inches(2.35), Inches(1.0))
    sh.fill.solid(); sh.fill.fore_color.rgb = LIGHT
    sh.line.color.rgb = BORDER; sh.line.width = Pt(0.5); sh.name = f"Stat5_{i}"
    tf2 = sh.text_frame; tf2.margin_top = Inches(0.08)
    p_v = tf2.paragraphs[0]; p_v.alignment = PP_ALIGN.CENTER
    r_v = p_v.add_run(); r_v.text = val
    set_font(r_v, size=26, bold=True, color=INDIGO)
    p_l = tf2.add_paragraph(); p_l.alignment = PP_ALIGN.CENTER
    r_l = p_l.add_run(); r_l.text = lbl
    set_font(r_l, size=8.5, color=GRAY)

oral(sl, '"La réservation déclenche une série d\'opérations : vérification SQL avec chevauchement de dates, attribution automatique de chambre physique, calcul du prix et notification Telegram avec boutons interactifs."')
footer(sl, 5)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DASHBOARD ADMIN
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_BLANK)
title_area(sl, "Admin React", "Dashboard administrateur",
           "Route /admin — protégée par JWT · 5 onglets")

cards6 = [
    ("Réservations", "Liste · Chambre attribuée · Changer statut"),
    ("Commandes",    "Room service · Statuts : en attente / préparation / livré"),
    ("Chambres",     "Vue temps réel : libre, occupée, arrivée, départ"),
    ("WiFi",         "Sessions RADIUS réelles + activité réseau démo"),
    ("Statistiques", "Revenus, taux d'occupation, activité du jour"),
    ("Chargement",   "Toutes les données en parallèle via Promise.all()"),
]
for i, (lbl, body) in enumerate(cards6):
    c = i % 3; r_i = i // 3
    card(sl, Inches(0.4)+c*Inches(4.17), Inches(1.55)+r_i*Inches(1.5),
         Inches(4.02), Inches(1.38), lbl, body)

card(sl, Inches(0.4), Inches(4.65), Inches(6.0), Inches(1.15),
     "Authentification Admin",
     "Mot de passe unique (.env) → JWT HS256 · 7 jours · sessionStorage")
card(sl, Inches(6.6), Inches(4.65), Inches(6.13), Inches(1.15),
     "Protection API",
     "AdminMiddleware vérifie le JWT à chaque requête → HTTP 401 si absent ou invalide")

oral(sl, '"L\'admin est une SPA React. Toutes les données sont chargées en parallèle. Chaque route API admin retourne HTTP 401 si le token est absent — il n\'y a aucun accès sans authentification."')
footer(sl, 6)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PORTAIL CAPTIF WIFI
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_BLANK)
title_area(sl, "WiFi — TP-Link EAP330", "Portail captif WiFi",
           "Authentification des clients par numéro de chambre")

flow7 = ["Téléphone → OP1_Wifi", "Portail EAP330", "Login 201/201",
         "RADIUS :18120", "FreeRADIUS Docker", "Access-Accept → Internet"]
tx = Inches(0.4)
bg7 = sl.shapes.add_shape(1, Inches(0.4), Inches(1.45), Inches(12.5), Inches(0.65))
bg7.fill.solid(); bg7.fill.fore_color.rgb = LIGHT
bg7.line.color.rgb = BORDER; bg7.line.width = Pt(0.5); bg7.name = "FlowBg7"
for i, step in enumerate(flow7):
    w = Inches(1.95)
    sh = sl.shapes.add_shape(1, tx, Inches(1.5), w, Inches(0.54))
    sh.fill.solid(); sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = BORDER; sh.line.width = Pt(0.5); sh.name = f"Flow7_{i}"
    tf2 = sh.text_frame; tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    r2 = tf2.paragraphs[0].add_run(); r2.text = step
    set_font(r2, size=7.5, color=DARK, bold=True)
    tx += w + Inches(0.12)
    if i < len(flow7)-1:
        arr = sl.shapes.add_textbox(tx-Inches(0.16), Inches(1.57), Inches(0.14), Pt(22))
        arr.text_frame.paragraphs[0].add_run().text = "→"
        set_font(arr.text_frame.paragraphs[0].runs[0], size=12, bold=True, color=INDIGO)

bullets(sl, [
    "Identifiant = Mot de passe = Numéro de chambre (ex: 201)",
    "Borne EAP330 en mode 'Local Web Portal' + External RADIUS Server",
    "Session autorisée 1 heure — puis re-authentification requise",
    "FreeRADIUS répond Access-Accept (OK) ou Access-Reject (refusé)",
    "RADIUS = protocole standard en hôtellerie, entreprises et universités",
], Inches(0.5), Inches(2.28), Inches(12.3), Inches(2.0))

# Badge démo
badge = sl.shapes.add_shape(1, Inches(0.5), Inches(4.5), Inches(6.5), Inches(0.4))
badge.fill.solid(); badge.fill.fore_color.rgb = rgb(0xEC,0xFD,0xF5)
badge.line.color.rgb = rgb(0x6E,0xE7,0xB7); badge.line.width = Pt(0.75)
badge.name = "DemoBadge"
tf_b = badge.text_frame; p_b = tf_b.paragraphs[0]
r_b = p_b.add_run(); r_b.text = "Démo live : connexion WiFi depuis un téléphone — chambre 201"
set_font(r_b, size=9.5, bold=True, color=GREEN)

oral(sl, '"La borne TP-Link EAP330 délègue l\'authentification à FreeRADIUS. Le client entre son numéro de chambre et accède à Internet. On a dû créer un relay UDP PowerShell car Docker Desktop Windows ne peut pas recevoir de paquets UDP du réseau local."')
footer(sl, 7)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FREERADIUS + ADMIN WIFI
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_BLANK)
title_area(sl, "FreeRADIUS + Remontée Admin", "Sessions WiFi dans l'admin")

bullets(sl, [
    "FreeRADIUS dans Docker : 10 chambres (101–603) configurées, secret partagé avec la borne",
    "Relay UDP PowerShell : écoute :18120 sur l'interface physique, retransmet vers FreeRADIUS :1812",
    "Script sync_radius_sessions.sh : parse logs FreeRADIUS → insère/met à jour MySQL (idempotent)",
    "Onglet WiFi admin : appareils connectés, chambres, adresses MAC, Access-Accept/Reject",
    "Activité réseau : section de démonstration clairement labelisée dans l'interface",
], Inches(0.5), Inches(1.35), Inches(12.3), Inches(2.2))

card(sl, Inches(0.5), Inches(3.75), Inches(6.0), Inches(1.15),
     "Données WiFi en base (réelles)",
     "4 sessions · 3 Access-Accept · 1 Access-Reject · Adresses MAC réelles")
card(sl, Inches(6.7), Inches(3.75), Inches(6.13), Inches(1.15),
     "Activité réseau (démonstration)",
     "10 lignes DNS · Chambres 101 et 201 · Marquées is_demo=1 en base")

badge2 = sl.shapes.add_shape(1, Inches(0.5), Inches(5.15), Inches(7.5), Inches(0.4))
badge2.fill.solid(); badge2.fill.fore_color.rgb = rgb(0xEC,0xFD,0xF5)
badge2.line.color.rgb = rgb(0x6E,0xE7,0xB7); badge2.line.width = Pt(0.75)
badge2.name = "DemoBadge2"
tf_b2 = badge2.text_frame; p_b2 = tf_b2.paragraphs[0]
r_b2 = p_b2.add_run(); r_b2.text = "Démo : bash sync_radius_sessions.sh → onglet WiFi mis à jour"
set_font(r_b2, size=9.5, bold=True, color=GREEN)

oral(sl, '"Un script bash lit les logs FreeRADIUS et les transforme en données SQL. L\'admin voit en temps réel quels appareils sont connectés et depuis quelle chambre. Le script peut être relancé sans créer de doublons (ON DUPLICATE KEY UPDATE)."')
footer(sl, 8)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — SÉCURITÉ & DIFFICULTÉS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_BLANK)
title_area(sl, "Sécurité & Problèmes résolus", "Difficultés et solutions")

cards9 = [
    ("JWT HS256",               "Tokens signés — impossible à forger sans le secret .env"),
    ("Mots de passe bcrypt",    "Jamais stockés en clair — hash irréversible en base"),
    ("Requêtes PDO préparées",  "Paramètres bindés — protection contre les injections SQL"),
    ("Ports restreints",        "phpMyAdmin + MySQL sur 127.0.0.1 uniquement"),
]
for i, (lbl, body) in enumerate(cards9):
    c = i % 2; r_i = i // 2
    card(sl, Inches(0.4)+c*Inches(6.3), Inches(1.55)+r_i*Inches(1.12),
         Inches(6.13), Inches(1.0), lbl, body)

# Séparateur
sep = sl.shapes.add_shape(1, Inches(0.4), Inches(3.88), Inches(12.5), Pt(1))
sep.fill.solid(); sep.fill.fore_color.rgb = BORDER; sep.line.fill.background()

tb_lbl = sl.shapes.add_textbox(Inches(0.4), Inches(3.98), Inches(12.3), Pt(14))
rl2 = tb_lbl.text_frame.paragraphs[0].add_run(); rl2.text = "PROBLÈMES RENCONTRÉS ET SOLUTIONS"
set_font(rl2, size=7.5, bold=True, color=GRAY)

# Problème/Solution dans un tableau natif
headers9 = ["Problème", "Cause", "Solution appliquée"]
rows9 = [
    ("Docker + UDP",
     "Docker Desktop Windows ne peut pas recevoir de paquets UDP depuis le réseau local (NAT)",
     "Relay UDP PowerShell : écoute 192.168.30.104:18120 → retransmet vers 127.0.0.1:1812"),
    ("Double réservation",
     "Deux clients pouvaient réserver la même chambre aux mêmes dates",
     "Requête SQL avec condition de chevauchement de dates + attribution automatique"),
    ("Boutons Telegram bloqués",
     "Connexion PDO singleton expirée (MySQL gone away) — answerCallbackQuery jamais appelé",
     "Fonction get_db() avec ping + Database::reset() + reconnexion automatique"),
]
add_table(sl, headers9, rows9, Inches(0.4), Inches(4.2), Inches(12.5), Inches(2.0))

oral(sl, '"La principale difficulté était Docker + RADIUS sur Windows. Le relay UDP PowerShell résout ce problème. On a aussi corrigé un bug : les boutons Telegram chargeaient à l\'infini car la connexion MySQL s\'était coupée."', t=Inches(6.33))
footer(sl, 9)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CONCLUSION
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_BLANK)
title_area(sl, "Bilan", "Conclusion et améliorations")

tb_ok = sl.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(12.3), Pt(14))
r_ok = tb_ok.text_frame.paragraphs[0].add_run(); r_ok.text = "CE QUI FONCTIONNE — DÉMONTRABLE EN LIVE"
set_font(r_ok, size=8, bold=True, color=GRAY)

bullets(sl, [
    "Réservation multi-nuits de bout en bout avec attribution automatique de chambre physique",
    "Room service : commande → Telegram → suivi statut dans l'admin",
    "Portail captif WiFi RADIUS — testé avec de vrais appareils sur le réseau",
    "Dashboard admin complet : réservations, commandes, chambres, WiFi, statistiques",
    "Infrastructure Docker stable — 8 services, IP statique, relay UDP persistant",
], Inches(0.5), Inches(1.62), Inches(12.3), Inches(1.85))

sep2 = sl.shapes.add_shape(1, Inches(0.4), Inches(3.62), Inches(12.5), Pt(1))
sep2.fill.solid(); sep2.fill.fore_color.rgb = BORDER; sep2.line.fill.background()

tb_amL = sl.shapes.add_textbox(Inches(0.5), Inches(3.72), Inches(12.3), Pt(14))
r_am = tb_amL.text_frame.paragraphs[0].add_run(); r_am.text = "AMÉLIORATIONS POSSIBLES"
set_font(r_am, size=8, bold=True, color=GRAY)

improv = [
    ("HTTPS",       "Certificat SSL (Let's Encrypt) pour chiffrer les échanges"),
    ("DNS réel",    "Pi-hole pour remplacer la démo par de vraies données réseau"),
    ("Téléphonie",  "Connecter des postes SIP physiques aux extensions Asterisk"),
    ("Monitoring",  "Alertes automatiques si un service Docker tombe"),
]
for i, (lbl, body) in enumerate(improv):
    card(sl, Inches(0.4)+i*Inches(3.1), Inches(4.0), Inches(2.95), Inches(1.25), lbl, body)

# Bas de slide sombre (fin du diaporama)
closing = sl.shapes.add_shape(1, 0, Inches(5.45), prs.slide_width, Inches(2.05))
closing.fill.solid(); closing.fill.fore_color.rgb = DARK
closing.line.fill.background(); closing.name = "ClosingBg"

tb_cl = sl.shapes.add_textbox(0, Inches(5.7), prs.slide_width, Inches(0.7))
p_cl = tb_cl.text_frame.paragraphs[0]; p_cl.alignment = PP_ALIGN.CENTER
r_cl = p_cl.add_run(); r_cl.text = "Maison Saclay — Hôtel 5 étoiles — Projet OP-1"
set_font(r_cl, size=16, bold=True, color=WHITE)

tb_cl2 = sl.shapes.add_textbox(0, Inches(6.3), prs.slide_width, Inches(0.5))
p_cl2 = tb_cl2.text_frame.paragraphs[0]; p_cl2.alignment = PP_ALIGN.CENTER
r_cl2 = p_cl2.add_run(); r_cl2.text = "Licence Pro MRT — 2025–2026"
set_font(r_cl2, size=11, color=GRAY)

oral(sl, '"Le projet est fonctionnel et démontrable dans son ensemble. Avec plus de temps : HTTPS, vrai DNS local, postes téléphoniques physiques et alertes monitoring."', t=Inches(7.05))
footer(sl, 10)

# ── Sauvegarde ───────────────────────────────────────────────────────────────
prs.save(out_path)
print("Fichier cree : " + out_path)
